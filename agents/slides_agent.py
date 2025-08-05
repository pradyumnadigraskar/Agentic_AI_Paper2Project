from utils.llm_utils import run_llm_prompt
from pptx import Presentation
import os

OUTPUT_FOLDER = 'output'

def create_slides(idea):
    prompt = (
        f"Create bullet points for 5 slides to present this idea: {idea}.\n"
        "Format like:\nSlide 1 - Title: ..., Bullets: ... ; ... ; ...\n"
        "Slide 2 - Title: ..., Bullets: ... ; ... ; ...\n"
    )
    return run_llm_prompt(prompt)

def generate_ppt(project_name, slides_text):
    prs = Presentation()
    slides_text = slides_text.strip()
    slide_chunks = slides_text.split("Slide ")

    for chunk in slide_chunks:
        if not chunk.strip():
            continue
        try:
            number_part, content_part = chunk.strip().split(" - ", 1)
        except ValueError:
            continue

        title = f"Slide {number_part.strip()}"
        if "Bullets:" in content_part:
            title_text, bullets_text = content_part.split("Bullets:", 1)
        else:
            title_text = content_part
            bullets_text = ""

        bullet_points = [b.strip() for b in bullets_text.split(";") if b.strip()]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text.strip()

        content_box = slide.placeholders[1]
        content_box.text = bullet_points[0] if bullet_points else ""

        for point in bullet_points[1:]:
            content_box.text_frame.add_paragraph().text = point

    os.makedirs(OUTPUT_FOLDER, exist_ok=True)
    ppt_path = os.path.join(OUTPUT_FOLDER, f"{project_name}_slides.pptx")
    prs.save(ppt_path)
    return ppt_path

# def save_slides_ppt(project_name, slides_text):
#     prs = Presentation()
#     slides_text = slides_text.strip()

#     # Split by 'Slide X -' entries
#     slide_chunks = slides_text.split("Slide ")
#     for chunk in slide_chunks:
#         if not chunk.strip():
#             continue
#         try:
#             number_part, content_part = chunk.strip().split(" - ", 1)
#         except ValueError:
#             continue

#         # Try parsing title and bullets
#         title = f"Slide {number_part.strip()}"
#         if "Bullets:" in content_part:
#             title_text, bullets_text = content_part.split("Bullets:", 1)
#         else:
#             title_text = content_part
#             bullets_text = ""

#         bullet_points = [b.strip() for b in bullets_text.split(";") if b.strip()]
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text = title_text.strip()

#         content_box = slide.placeholders[1]
#         content_box.text = bullet_points[0] if bullet_points else ""

#         for point in bullet_points[1:]:
#             content_box.text_frame.add_paragraph().text = point

#     os.makedirs(OUTPUT_FOLDER, exist_ok=True)
#     ppt_path = os.path.join(OUTPUT_FOLDER, f"{project_name}_slides.pptx")
#     prs.save(ppt_path)
#     return ppt_path
