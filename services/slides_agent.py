from pptx import Presentation
import os

def create_slides(idea, llm_provider):
    print("LOG: Generating Slide Content...")
    prompt = f"Create a 5-slide presentation outline for: {idea}. Format: Slide X: Title | Bullets: point1; point2"
    return llm_provider.generate(prompt)

def generate_ppt(project_name, slides_text):
    print("LOG: Creating PPTX file...")
    prs = Presentation()
    raw_slides = slides_text.split("Slide")
    for slide_data in raw_slides:
        if len(slide_data.strip()) < 5: continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = slide_data.strip().split('\n')
        slide.shapes.title.text = lines[0].split('|')[0].strip()
        tf = slide.placeholders[1].text_frame
        for line in lines[1:]:
            clean = line.strip().lstrip("-*•|")
            if clean: tf.add_paragraph().text = clean

    file_name = f"{project_name}_slides.pptx"
    prs.save(os.path.join("output", file_name))
    return file_name